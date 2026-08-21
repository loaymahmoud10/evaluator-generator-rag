"""Load document-based knowledge sources into LangChain Documents."""

from __future__ import annotations

from pathlib import Path

from docx import Document as DocxDocument
from langchain_community.document_loaders import PyMuPDFLoader
from langchain_core.documents import Document
from pptx import Presentation
from app.ingestion.base import BaseLoader

# Common source-code extensions treated as textual knowledge.
CODE_EXTENSIONS = {
    ".py",
    ".js",
    ".ts",
    ".java",
    ".c",
    ".cpp",
    ".h",
    ".cs",
    ".go",
    ".rs",
    ".rb",
    ".php",
    ".sql",
    ".sh",
    ".html",
    ".css",
    ".json",
    ".yaml",
    ".yml",
    ".xml",
    ".md",
}


class PDFLoader(BaseLoader):
    """Load PDF files while preserving source metadata."""

    def load(self, source: str | Path) -> list[Document]:
        """Load a PDF and return its pages as LangChain Documents."""
        path = Path(source)

        if not path.exists():
            raise FileNotFoundError(f"PDF file not found: {path}")

        if path.suffix.lower() != ".pdf":
            raise ValueError(f"Expected a PDF file, got: {path.suffix}")

        loader = PyMuPDFLoader(str(path))
        documents = loader.load()

        source_id = path.resolve().as_posix()

        for document in documents:
            page_number = document.metadata.get("page")

            document.metadata.update(
                {
                    "source_id": source_id,
                    "source_type": "pdf",
                    "source_name": path.name,
                    "location": (
                        f"page {page_number + 1}"
                        if isinstance(page_number, int)
                        else "unknown"
                    ),
                }
            )

        return documents


class DOCXLoader(BaseLoader):
    """Load DOCX files while preserving source metadata."""

    def load(self, source: str | Path) -> list[Document]:
        """Load a DOCX file and return it as a LangChain Document."""
        path = Path(source)

        if not path.exists():
            raise FileNotFoundError(f"DOCX file not found: {path}")

        if path.suffix.lower() != ".docx":
            raise ValueError(f"Expected a DOCX file, got: {path.suffix}")

        docx = DocxDocument(str(path))

        text = "\n".join(
            paragraph.text
            for paragraph in docx.paragraphs
            if paragraph.text.strip()
        )

        document = Document(
            page_content=text,
            metadata={
                "source_id": path.resolve().as_posix(),
                "source_type": "docx",
                "source_name": path.name,
                "location": "document",
            },
        )

        return [document]

class TXTLoader(BaseLoader):
    """Load plain-text files while preserving source metadata."""

    def load(self, source: str | Path) -> list[Document]:
        """Load a text file and return it as a LangChain Document."""
        path = Path(source)

        if not path.exists():
            raise FileNotFoundError(f"TXT file not found: {path}")

        if path.suffix.lower() != ".txt":
            raise ValueError(f"Expected a TXT file, got: {path.suffix}")

        text = path.read_text(encoding="utf-8")

        document = Document(
            page_content=text,
            metadata={
                "source_id": path.resolve().as_posix(),
                "source_type": "txt",
                "source_name": path.name,
                "location": "document",
            },
        )

        return [document]

class CodeLoader(BaseLoader):
    """Load source-code files as textual knowledge."""

    def load(self, source: str | Path) -> list[Document]:
        """Load a source-code file and return it as a LangChain Document."""
        path = Path(source)

        if not path.exists():
            raise FileNotFoundError(f"Code file not found: {path}")

        if not path.is_file():
            raise ValueError(f"Expected a file, got: {path}")

        text = path.read_text(encoding="utf-8")

        document = Document(
            page_content=text,
            metadata={
                "source_id": path.resolve().as_posix(),
                "source_type": "code",
                "source_name": path.name,
                "location": "document",
            },
        )

        return [document]

class PPTXLoader(BaseLoader):
    """Load PPTX presentations while preserving slide-level source metadata."""

    def load(self, source: str | Path) -> list[Document]:
        """Load a PPTX file and return one LangChain Document per slide."""
        path = Path(source)

        if not path.exists():
            raise FileNotFoundError(f"PPTX file not found: {path}")

        if path.suffix.lower() != ".pptx":
            raise ValueError(f"Expected a PPTX file, got: {path.suffix}")

        presentation = Presentation(str(path))
        documents: list[Document] = []

        source_id = path.resolve().as_posix()

        for slide_number, slide in enumerate(presentation.slides, start=1):
            text_parts: list[str] = []

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text.strip())

            slide_text = "\n".join(text_parts)

            documents.append(
                Document(
                    page_content=slide_text,
                    metadata={
                        "source_id": source_id,
                        "source_type": "pptx",
                        "source_name": path.name,
                        "location": f"slide {slide_number}",
                    },
                )
            )

        return documents